import os
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.text_splitter import CharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores.faiss import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain_openai import OpenAI
from dotenv import load_dotenv
load_dotenv()
os.environ["OPENAI_API_KEY"] = os.getenv("OPENAI_API_KEY")

text_splitter=CharacterTextSplitter(
    separator="\n",
    chunk_size=1000,
    chunk_overlap=200,
    length_function = len)

embeddings = OpenAIEmbeddings()
chain = load_qa_chain(OpenAI(max_tokens=512, temperature=0))
raw_text = ""
doc_search = FAISS

def create_doc_search(text):
    texts=text_splitter.split_text(text)
    doc_search= FAISS.from_texts(texts, embeddings)
    return doc_search

def process_files(files):
    global raw_text
    for file in files:
        if(file.type == "application/pdf"):
            pdfReader = PdfReader(file)
            for page in pdfReader.pages:
                raw_text += page.extract_text()
        elif(file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        raw_text += shape.text  
    global doc_search
    doc_search=create_doc_search(raw_text)
                    
def answer_question(prompt):
    
    docs=doc_search.similarity_search(query=prompt)
    
    chain.llm_chain.prompt.template="You are an intelligent agent who specialises in design thinking, you do not have any knowledge other than the context provided below. ONLY use the context given below to answer the question at the end. Include examples wherever possible. Strictly use bullet points. Only use the information provided in the context. ONLY answer the question asked, do not give miscellaneous content. If you don't know the answer, just say that you don't know, don't try to make it up. context: {context}  question: {question}" #a prompt should include {context} and {question} placeholders, include stuff like how many words answer for 5 marks and or 10 marks and for 7-8 marks. tell it to give examples when explaining concepts
    
    answer=chain.run(input_documents=docs, question=prompt)
    return answer